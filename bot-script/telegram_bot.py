from __future__ import annotations

from collections import defaultdict

from dotenv import load_dotenv
import base64
import signal
import asyncio
import json
import logging
import os
import re
import io
import threading
import time
import traceback
from datetime import datetime, timedelta
from signal import SIGINT, SIGTERM, SIGABRT

from openai_helper import default_max_tokens

import schedule
import telegram
import concurrent.futures
import payment
from uuid import uuid4

from yookassa import Configuration, Payment

from telegram import BotCommandScopeAllGroupChats, Update, constants
from telegram import InlineKeyboardMarkup, InlineKeyboardButton, InlineQueryResultArticle
from telegram import InputTextMessageContent, BotCommand
from telegram.error import RetryAfter, TimedOut
from telegram.ext import ApplicationBuilder, CommandHandler, MessageHandler, \
    filters, InlineQueryHandler, CallbackQueryHandler, Application, ContextTypes, CallbackContext, PollAnswerHandler

from pydub import AudioSegment

from utils import is_group_chat, get_thread_id, message_text, wrap_with_indicator, split_into_chunks, \
    edit_message_with_retry, get_stream_cutoff_values, is_allowed, is_admin, \
    get_reply_to_message_id, error_handler
from openai_helper import OpenAIHelper, localized_text
# from usage_tracker import UsageTracker

from aiogram import Bot, Dispatcher, executor, types
from db import Database
from db_analytics import *

import os
import sys

from asgiref.sync import sync_to_async
from django.forms import model_to_dict

# Получаем путь к текущему скрипту
script_path = os.path.abspath(__file__)

# Получаем путь к директории, содержащей текущий скрипт
script_dir = os.path.dirname(script_path)

# Получаем путь к корневой директории проекта (по одному уровню выше)
project_root = os.path.dirname(script_dir)

# Добавляем путь к корневой директории в переменную окружения PYTHONPATH
sys.path.insert(0, project_root)

# Теперь можно импортировать модели из bot.models

import PyPDF2
import docx
from pptx import Presentation
import openpyxl
import io

class ChatGPTTelegramBot:
    """
    Class representing a ChatGPT Telegram Bot.
    """

    def __init__(self, config: dict, openai: OpenAIHelper):
        """
        Initializes the bot with the given configuration and GPT bot object.
        :param config: A dictionary containing the bot configuration
        :param openai: OpenAIHelper object
        """
        self.interrupt_flag = False
        self.notif_run = False
        self.tasks = []
        self.db = Database()

        self.db_analytics_for_sessions = DBanalytics_for_sub_stat()
        self.db_statistic_by_day = DBstatistics_by_day()
        self.db_admin_stats = DBAdminStats()


        self.config = config
        self.openai = openai

        self.status = {}
        self.prompts: dict[int: list] = {}
        bot_language = self.config['bot_language']
        self.commands = [
            BotCommand(command='help', description='Помощь/описание'),
            BotCommand(command='reset', description='Сбросить историю'),
            BotCommand(command='buy', description='Купить подписку'),
            BotCommand(command='stats', description='Моя Статистика'),
            BotCommand(command='resend', description='Переслать последний запрос'),
            BotCommand(command='save', description='Закрепить выбранное сообщение'),
            BotCommand(command='delete', description='Удалить из контекста выбранное сообщение'),
            BotCommand(command='model', description='Изменить модель'),
            BotCommand(command='imagine', description='Сгенерировать изображение'),
            BotCommand(command='quality', description='Изменить качество изображения'),



        ]
        self.commands.append(BotCommand(command='role', description='Изменить роль  PRO'))
        self.commands.append(BotCommand(command='temperature', description='Изменить креативность  PRO'))

        self.group_commands = [BotCommand(
            command='chat', description=localized_text('chat_description', bot_language)
        )] + self.commands
        self.disallowed_message = localized_text('disallowed', bot_language)
        self.budget_limit_message = localized_text('budget_limit', bot_language)
        self.usage = {}
        self.last_message = {}
        self.quality_list = {}
        self.inline_queries_cache = {}
        self.bot = Bot(token=self.config['token'])
        self.dispatcher = Dispatcher(bot=self.bot, loop=asyncio.get_event_loop())
        self.poll_results = defaultdict(int)

    async def test_poll(self, update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:

        if not is_admin(self.config, update.message.from_user.id):
            return


        admin_users = await self.db.get_admin_users()


        for user in admin_users:
            await self.bot.send_message(
                chat_id=user,
                text='Друзья, мы хотим сделать бот лучше! И нам нужна ваша помощь. Расскажите:',
            )


            message = await self.bot.send_poll(
                chat_id=user,
                question='Почему вы не готовы купить подписку?',
                options=['Пока не надо', 'Дорого — у других выгоднее', 'Хочу больше моделей, функций'],
                is_anonymous=False,
                allows_multiple_answers=False,
            )




            payload = {
                message.poll.id: message.message_id,
                user: message.poll.id,


            }
            context.bot_data.update(payload)

    async def poll(self, update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:

        if not is_admin(self.config, update.message.from_user.id):
            return


        users = await self.db.get_all_users()


        for user in users:
            err = 0
            success = 0
            try:
                await self.bot.send_message(
                    chat_id=user,
                    text='Друзья, мы хотим сделать бот лучше! И нам нужна ваша помощь. Расскажите:',
                )


                message = await self.bot.send_poll(
                    chat_id=user,
                    question='Почему вы не готовы купить подписку?',
                    options=['Пока не надо', 'Дорого — у других выгоднее', 'Хочу больше моделей, функций'],
                    is_anonymous=False,
                    allows_multiple_answers=True,
                )





                payload = {
                    message.poll.id: message.message_id,
                    user: message.poll.id,


                }
                context.bot_data.update(payload)
                success += 1
            except Exception as e:
                err += 1


        logging.info(f'Errors: {err}')
        logging.info(f'Success: {success}')




    async def receive_poll_answer(self, update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
        """Summarize a users poll vote"""

        try:
            logging.info('Poll answer: %s', update.poll_answer)
            answer = update.poll_answer
            selected_options = answer.option_ids
            answer = ''
            for option_id in selected_options:
                self.poll_results[option_id] += 1
                answer += f'{option_id}' + ' '
            logging.info('Poll results: %s', self.poll_results)
            await self.db.add_poll_answer(user_id=update.poll_answer.user.id, answer=answer)

        except Exception as e:
            logging.error(f'Error adding poll answer: {e}')
            pass




    async def show_results(self, update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:
        # check admin
        if not is_admin(self.config, update.message.from_user.id):
            return
        result_message = ""
        for option in self.poll_results:
            result_message += f"{option}: {self.poll_results[option]}\n"

        if result_message:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text=result_message,
            )




    async def send_info_message(self, users, text):

        try:
            for user in users:
                await self.bot.send_message(
                    chat_id=user.user_id,
                    text=text,
                )
        except Exception as e:
            logging.error(f'Error sending info message: {e}')
            pass

    async def add_client(self, update: Update, _: ContextTypes.DEFAULT_TYPE, user_id, client_id):
        import requests
        import json
        from dotenv import load_dotenv
        load_dotenv()


        url = " https://api-metrika.yandex.net/cdp/api/v1/counter/94971306/data/contacts?merge_mode=APPEND"
        token = os.environ.get('METRICS_BOT_TOKEN')

        headers = {
            'Authorization': token,
            "Content-Type": "application/json",
        }

        import pytz

        # Создаем объект datetime
        dt = datetime.now() - timedelta(seconds=60)
        dt = dt.strftime("%Y-%m-%d %H:%M:%S")


        data = {
    "contacts":
    [
        {
            "uniq_id" : str(user_id),

            "create_date_time" : dt,

            "client_ids": [str(client_id) ],


        }
    ]
}

        response = requests.post(url, headers=headers, data=json.dumps(data))

        print(response.json())





    async def add_order(self, user_id, revenue, cost, order_id, product):
        import requests
        import json

        load_dotenv()
        token = os.environ.get('METRICS_BOT_TOKEN')
        url = "https://api-metrika.yandex.net/cdp/api/v1/counter/94971306/data/orders?merge_mode=APPEND"

        headers = {
            'Authorization': token,
            'Content-Type': 'application/json',
        }
        # minus 60 sec
        date = datetime.now() - timedelta(seconds=60)
        date = date.strftime("%Y-%m-%d %H:%M:%S")

        data = {
            "orders": [
        {
            "id" : str(order_id),
            "client_uniq_id" : str(user_id),
            "client_type" : "CONTACT",
            "create_date_time": date,
            "revenue" : revenue,
            "order_status" : "1",
            "cost" : cost,

        }
        ]
        }

        response = requests.post(url, headers=headers, data=json.dumps(data))

        print(response.status_code)
        print(response.json())


    async def add_offline(self, user_id, target):
        client_id = await self.db.get_client_id(user_id)
        if not client_id:
            return
        
        if not await self.db.check_offline_conversions_settings_count(target):
            return
        
        if await self.db.get_offline_conversions_of_user(user_id, target):
            return

        import requests
        import csv
        import json
        id_type = "CLIENT_ID"
        load_dotenv()
        counter = 94971306
        token = os.environ.get('METRICS_BOT_TOKEN')
        url = "https://api-metrika.yandex.net/management/v1/counter/{}/offline_conversions/upload?client_id_type={}".format(counter, id_type)
        headers = {
        "Authorization": token
        
        }
        # Вычисляем текущее время с учетом задержки (на 1 минуту меньше текущего времени)
        # Дата и время конверсии в формате Unix Time Stamp
        date = datetime.now() - timedelta(seconds = 3 )
        
        date = int(date.timestamp())
        date = str(date)
    

        logging.info(f'{client_id} {target} {date}')


     

        data = {
            'ClientId' : client_id,
            'Target' : target,
            'DateTime' : date,
        }
        # Создаем CSV файл на лету с расделитеем ,
        output = io.StringIO()
        csv_writer = csv.writer(output, delimiter=',')
        csv_writer.writerow(data.keys())  # Пишем заголовки
        csv_writer.writerow(data.values())  # Пишем данные

        # Перемещаем курсор на начало файла
        output.seek(0)



    
        # Отправляем запрос
        files = {'file': ('offline-conversions.csv', output, 'text/csv')}


        

        response = requests.post(url, headers=headers, files=files)

        # Закрываем StringIO объект
        output.close()

        logging.info(response.status_code)
        logging.info(response.json())



        response.raise_for_status()

        if response.status_code == 200:
            
            await self.db.add_offline_conversions(user_id, target)
            await self.db.add_offline_conversions_settings_count(target)
        
        else:
            logging.error(f"Error adding offline conversions: {response.status_code}")
            logging.error(f"Error adding offline conversions: {response.json()}")

        
        

        


    async def orders(self, update: Update, _: ContextTypes.DEFAULT_TYPE):
        import requests
        import json
        load_dotenv()
        token = os.environ.get('METRICS_BOT_TOKEN')


        headers = {
            'Authorization': token,
            'Content-Type': 'application/json',
        }


        url = "https://api-metrika.yandex.net/cdp/api/v1/counter/94971306/schema/order_statuses"



        data = {
            "order_statuses": [
            {
                "id": "1",
                "type": "PAID"
            }]
        }

        response = requests.get(url, headers=headers, data=json.dumps(data))

        print(response.json())




    

    async def cancel(self, update: Update, _: ContextTypes.DEFAULT_TYPE) -> None:
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text="Отменено",
        )
        user_id = update.message.from_user.id
        self.status[user_id] = 'prompt'

    async def help(self, update: Update, _: ContextTypes.DEFAULT_TYPE) -> None:
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            parse_mode='HTML',
            text='''🔻Если возникли вопросы, столкнулись с ошибкой напишите нам brainshtorm@gmail.com
            
<b>📈 Как считаются токены</b>

<blockquote><b>Одна тыс.</b> токенов примерно равна 1.5 стр. А4. 
❗️Но❗️
Потраченные токены зависят от длины <b>вопроса и ответа</b> GPT. Всё вместе называется <b>история</b> или же контекст. Это третий параметр, который <b>влияет на потраченные токены.</b> И тратит он больше всех.

<b>Поэтому не забывайте сбрасывать историю командой /reset.</b> Так вы прочистите мозги нейросети, а еще сэкономите токены.

Под каждым ответом GPT написано количество токенов, которые находятся в контексте, по-другому – истории чата.
</blockquote>

<b>📚 Как считается история</b>

<blockquote>Вопрос + Ответ = История 1 запроса; История 1 запроса + Вопрос 2 + Ответ 2 = История 2 запроса.

<b>Короче говоря, каждый следующий запрос к нейросети включает в себя все предыдущие вопросы и ответы.</b>

Не ругайтесь на нас. Это придумали не мы. Убрать это нельзя. Мы же хотим предупредить вас об этом. Поэтому чистите историю чаще /reset :)
</blockquote>
<b>⚙️ Команда model</b>

<blockquote><b>В подписках с GPT-4 включена и GPT-4 mini.</b> С помощью команды /model можно переключаться между моделями. <b>Но расход токенов при «GPT-4-mini» уменьшается в 5 раз.</b> Получается, что 40 000 токенов «GPT-4» фактически 200 000 в GPT-4-mini
</blockquote>
<b>🎭 Команда role</b>

<blockquote>После ввода команды /role вы пишете <b>условия, которые нейросеь должна соблюдать.</b> Например, если нужны краткие ответы без пояснений, можно попросить нейросеть отвечать только "да" или "нет".
</blockquote>
<b>💡 Команда temperature</b>

<blockquote><b>Команда /temperature для того, чтобы регулировать креативность от 0 до 1.25.</b> Чем меньше температура, тем чаще GPT повторяется, но уменьшается шанс ошибки. Чем выше, тем креативнее и безумнее нейросеть. Начальное, самое стабильное значение 1
</blockquote>
<b>🧷 Команда save</b>

<blockquote><b>Команда /save – это просто закреплённые сообщения.</b> Например, чтобы не листать весь диалог с ответами на экзамене теперь можно сделать «точки» навигации. И вот как:

- Свйпните влево сообщение, которое хотите закрепить. На ПК – кликнуть правой кнопкой по сообщению и нажать в списке «ответить»
- Введите команду /save и отправьте её
- Всё, готово :)

Закрепить можно любые сообщения: свои и GPT. Количество не ограничено.
</blockquote>
<b>❌ Команда delete</b>

<blockquote><b>Эта команда удаляет выбранное сообщение из истории.</b> Delete нужен для того, чтобы тонко настраивать контекст чата. Допустим, нейросеть ответила не так, как вы хотели. И чтобы GPT не использовал плохой ответ, как источник информации для следующих ответов – просто удалите его:

- Свайпните влево сообщение, которое хотите закрепить. На ПК два ЛКМ по нему
- Введите команду /delete и отправьте
- Готово, сообщение удалено из контекста
</blockquote>
<b>📊 Анализ фото</b>

<blockquote>Доступен только тем, у кого есть подписка с GPT-4. Просто прикрепляете фото в чат, задаёте вопрос и всё. Готово. <b>Но помните, что одна фотография весит 1500 токено</b>
</blockquote>
<b>📸 Команды /imagine и /quality</b>

<blockquote>Нейросеть, которая генерирует фото - это DALL-E 3. Её создатели компания OpenAI,  это они разработали ChatGPT.  В её картинках меньше артефактов, больше деталей, но с её помощью пока нельзя обрабатывать фото. Чтобы сгенерировать изображение введите команду /imagine. После отдельным сообщением введите запрос, чт именно вы хотите получить на фото. Команда /quality - это выбор качества между Standard и HD.
</blockquote>
Подробнее на сайте: brainstormai.ru
''',

        )

    async def start(self, update: Update, context: ContextTypes.DEFAULT_TYPE) -> None:


        user_id = update.message.from_user.id

        logging.info(f'User {update.message.from_user.name} (id: {user_id}) started the bot')

        

        self.status[user_id] = 'prompt'





        if not await self.db.user_exists(user_id):

            try:
                await self.db.add_user(user_id)
            except Exception as e:
                logging.error(f'Error adding user {user_id} to the database: {e}')
            await self.calc_end_time(user_id)
            sub_id = await self.db.get_sub_type(user_id)
            try:


                await self.db_analytics_for_sessions.new_sub_stats(user_id, sub_id)
            except Exception as e:

                pass

            try:
                arg = update.message.text[7:]

                arg = arg.split('_')
                print(arg, user_id)
                # если в списке нет 5 элементов, запонить его None
                if len(arg) < 3:
                    for i in range(3 - len(arg)):
                        arg.append('')


                await self.db.set_utm(user_id, arg[0], arg[1], arg[2])




            except Exception as e:

                pass
                
            try:

                await self.add_offline(user_id, 'commandstart')
            except Exception as e:
                logging.error(e)
                logging.error('Error adding offline on start')
            pass

            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                parse_mode='HTML',

                text='''Добро пожаловать!

🆓 Активная подписка: Пробный период

⏬ Вам дотупно ⏬

✅ Дней: 3

✅ Модель: GPT-4o-mini

✅ Токенов: 3000 в день

<b>Важно</b>🔻

Потраченные токены зависят от длины вопроса и ответа GPT. Всё вместе называется стория или контекст. Это третий параметр, который влияет на потраченные токены. И тратит он больше всех.

Поэтому не забывайте сбрасывать историю командой /reset. Так вы «прочистите мозги» нейросети, а еще сэкономите токены.

Под каждым ответом GPT написано количество токенов, которые находятся в контексте, по-другому истории чата.

Подробнее /help
''',
            )

            



            return
        else:
            await self.db.set_unblocked_user(user_id)



    async def save(self, update:Update, context: ContextTypes.DEFAULT_TYPE):
        if await self.is_active(update, context, update.message.from_user.id) == False:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка закончилась, купите подписку',
            )
            return
        user_id = update.message.from_user.id

        # pin last bot message
        if update.message.reply_to_message:
            await update.message.reply_to_message.pin()
        else:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Вы не выбрали сообщение. Для выбора свайпните его влево. На ПК – 2 раза кликнуть по сообщению',
            )

    async def delete(self, update:Update, context: ContextTypes.DEFAULT_TYPE):
        if await self.is_active(update, context, update.message.from_user.id) == False:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка закончилась, купите подписку',
            )
            return
        user_id = update.message.from_user.id

        # pin last bot message
        if update.message.reply_to_message:
            # убрать 27 симовлов с конца
            text_message = update.message.reply_to_message.text

             # list od dicts
            deleted= False

            if self.openai.conversations.get(user_id):
                for item in self.openai.conversations[user_id]:


                    if item['role'] == 'user' and item['content'] == text_message:
                        self.openai.conversations[user_id].remove(item)
                        deleted = True
                        break
                    elif item['role'] == 'assistant':

                        short_text_message = text_message[0:len(text_message) - 28]

                        if item['content'] in short_text_message:
                            self.openai.conversations[user_id].remove(item)
                            deleted = True
                            break

                await update.message.reply_to_message.reply_text('Сообщение удалено')

            if deleted == False:
                await update.message.reply_to_message.reply_text('Сообщение не найдено')

        else:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Вы не выбрали сообщение. Для выбора свайпните его влево. На ПК – 2 раза кликнуть по сообщению',
            )
    def get_quality(self, user_id):
        if self.quality_list.get(user_id) == None:
            return {'quality': "standard", 'size': "1024x1024", 'flag' : True, 'tokens': 5000, 'price': 0.040}
        if self.quality_list.get(user_id,'st-1') == 'st-1':
            return {'quality': "standard", 'size': "1024x1024", 'tokens': 5000, 'price': 0.040}
        elif self.quality_list.get(user_id,'st-1') == 'st-2':
            return {'quality': "standard", 'size': "1024x1792", 'tokens': 2500, 'price': 0.080}
        elif self.quality_list.get(user_id,'st-1') == 'st-3':
            return {'quality': "standard", 'size': "1792x1024", 'tokens': 2500, 'price': 0.080}
        elif self.quality_list.get(user_id,'st-1') == 'hd-1':
            return {'quality': "hd", 'size': "1024x1024", 'tokens': 10000, 'price': 0.080}
        elif self.quality_list.get(user_id,'st-1') == 'hd-2':
            return {'quality': "hd", 'size': "1024x1792", 'tokens': 5000, 'price': 0.120}
        elif self.quality_list.get(user_id,'st-1') == 'hd-3':
            return {'quality': "hd", 'size': "1792x1024", 'tokens': 5000, 'price': 0.120}



    async def quality(self, update:Update, context: ContextTypes.DEFAULT_TYPE):
        if await self.is_active(update, context, update.message.from_user.id) == False:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка закончилась, купите подписку',
            )
            return

        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Выберите качество',
            reply_markup=InlineKeyboardMarkup([
                [InlineKeyboardButton('Standard 1024x1024 - 5000 токенов', callback_data='st-1')],
                # [InlineKeyboardButton('Standard 1024x1792 - 1000 токенов', callback_data='st-2')],
                # [InlineKeyboardButton('Standard 1792x1024 - 1000 токенов', callback_data='st-3')],
                [InlineKeyboardButton('HD 1024x1024 - 10000 токенов', callback_data='hd-1')],
                # [InlineKeyboardButton('HD 1024x1792 - 1000 токенов', callback_data='hd-2')],
                # [InlineKeyboardButton('HD 1792x1024 - 1000 токенов', callback_data='hd-3')],
            ])
        )

    async def fluxpro(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        user_id = update.message.from_user.id

        self.status[user_id] = 'fluxpro'
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text=self.status[user_id],
        )
        return

    async def fluxdev(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        user_id = update.message.from_user.id

        self.status[user_id] = 'fluxdev'
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text=self.status[user_id],
        )
        return
    async def imagine(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        if await self.is_active(update, context, update.message.from_user.id) == False:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка закончилась, купите подписку',
            )
            return

        sub_id = await self.db.get_sub_type(update.message.from_user.id)
        if await self.db.get_gen_im(sub_id) == False:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка не позволяет генерировать изображения',
            )
            return


        user_id = update.message.from_user.id





        res = self.get_quality(user_id)
        flag = res['flag'] if 'flag' in res else False
        quality = res['quality'] if 'quality' in res else 'standard'
        size = res['size'] if 'size' in res else '1024x1024'
        tokens = res['tokens'] if 'tokens' in res else 1000
        price = res['price'] if 'price' in res else 0.040




        if not await self.is_input_in_tokens(update, context, update.message.from_user.id, tokens):
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='У вас недостаточно токенов, выберите другое качество или купите подписку',
            )
            return
        if flag:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='На данный момент у вас установлено качество ' + quality + ' и размер ' + size + '. Чтобы изменить качество и размер, введите /quality',
            )
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Введите запрос',
        )

        self.status[user_id] = 'image'




    async def stats(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        # send message hello
        if await self.is_active(update, context, update.message.from_user.id) == False:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка закончилась, купите подписку',
            )
            return



        remain_tokens = await self.db.get_max_tokens(update.message.from_user.id) - await  self.db.get_used_tokens(
            update.message.from_user.id)

        date = str(await self.db.get_end_time(update.message.from_user.id))[0:10]
        date = date[8:10] + '.' + date[5:7] + '.' + date[0:4]

        await update.message.reply_text(
                text='Осталось: ' + str(
                    remain_tokens) + ' токенов' + '\n' + 'Подписка: ' + await self.db.get_sub_name_from_user(
                    update.message.from_user.id) + '\n' + 'Закончится: ' +
                     date
            )


    async def resend(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        """
        Resend the last request
        """
        chat_id = update.effective_chat.id
        if chat_id not in self.last_message:
            logging.warning(f'User {update.message.from_user.name} (id: {update.message.from_user.id})'
                            f' does not have anything to resend')
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text=localized_text('resend_failed', self.config['bot_language'])
            )
            return

        # Update message text, clear self.last_message and send the request to prompt
        logging.info(f'Resending the last prompt from user: {update.message.from_user.name} '
                     f'(id: {update.message.from_user.id})')
        with update.message._unfrozen() as message:
            message.text = self.last_message.pop(chat_id)

        await self.prompt(update=update, context=context)

    async def reset(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        """
        Resets the conversation.
        """

        logging.info(f'Resetting the conversation for user {update.message.from_user.name} '
                     f'(id: {update.message.from_user.id})...')

        chat_id = update.effective_chat.id
        reset_content = message_text(update.message)
        self.openai.reset_chat_history(chat_id=chat_id, content=reset_content)
        await update.effective_message.reply_text(
            message_thread_id=get_thread_id(update),
            text='История чата сброшена',
        )

    def change_model_of_sub(self, sub_name,curent_model):



        if sub_name in ['Multi Light', 'Multi PRO', 'Multi Standart', 'ultimate admin']:
            if curent_model in ['gpt-3.5']:
                return 'llama-3-70'
            elif curent_model in ['gpt-4']:
                return 'gpt-3.5'
            elif curent_model in ['llama-3-70']:
                return 'gpt-4'
        elif sub_name in ['Multi Mini']:
            if curent_model in ['gpt-3.5']:
                return 'llama-3-70'
            elif curent_model in ['llama-3-70']:
                return 'gpt-3.5'






    async def model(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        user_id = update.message.from_user.id
        self.status[user_id] = 'prompt'
        if await self.is_active(update, context, update.message.from_user.id) == False:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка акончилась, купите подписку',
            )
            return
        sub_id = await self.db.get_sub_type(update.message.from_user.id)

        if await self.db.get_sub_multimodel(sub_id):

            model = await self.db.get_user_model(user_id)

            new_model = await self.db.set_next_model(sub_id, user_id)

            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Модель изменена на ' + new_model,
            )

        else:
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка не позволяет менять модель, подробнее в /buy',
            )




    async def send_to_admin(self, text):
        try:
            admins = await self.db.get_admin_users()
            for admin in admins:
                await self.bot.send_message(chat_id=admin,
                                            text=text)
        except:
            pass




    async def send_to_all(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        if not await  self.db.is_admin(update.message.from_user.id):
            return
        self.status[update.message.from_user.id] = 'admin_message'
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Введите сообщение',
        )
    async def send_to_act(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        if not await  self.db.is_admin(update.message.from_user.id):
            return
        self.status[update.message.from_user.id] = 'admin_message_to_act'
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Введите сообщение',
        )

    async def send_reminder(self):



            try:
                users = await self.db.get_all_inactive_users()
                count = len(users)
                count_error = 0
                error_messages = []

                for user in users:
                    try:
                        await self.bot.send_message(chat_id=user,
                                                    text='Привет! Ты давно не заходил к нам. Наш бот всегда готов помочь тебе! Надеемся увидеть тебя снова')

                    except Exception as e:
                        count_error += 1
                        error_messages.append(str(e))
                        await self.db.set_blocked_user(user)
                        pass

                unique_error_messages = (set(error_messages))
                await self.send_to_admin('Отправили напоминание' + '\n' + 'Количество пользователей подошло: ' + str(count) + '\n'+
                                                    'Количество пользователей с ошибкой: ' + str(count_error) + '\n' + 'Уникальные ошибки: ' + str(unique_error_messages))

            except Exception as e:
                print(traceback.format_exc())
                await self.send_to_admin( 'error in send reminder' + '\n' + str(e))





    async def send_notif(self):



        async def job():
            await self.send_to_admin('start send notif')
            await self.send_notifications()
            await self.send_to_admin('start send reminders')
            await self.send_reminder()


        def run_job():
            asyncio.create_task(job())



        schedule.every().day.at('00:00').do(run_job)

        try:
            while True:
                schedule.run_pending()
                await asyncio.sleep(60)
        except KeyboardInterrupt:
            print('exit')

    # def interrupt(self):
    #     self.interrupt_flag = True
    #     for task in self.tasks:
    #         task.cancel()
    #     asyncio.gather(*self.tasks, return_exceptions=True)

    async def send_notifications(self):
            try:
                admin = await self.db.get_admin_users()
                users_for_inactive = await self.db.set_inactive_auto()
                count_users_for_inactive = len(users_for_inactive)
                count_error_users_for_inactive = 0
                if count_users_for_inactive > 0:
                    for user in users_for_inactive:
                        try:
                            await self.db_analytics_for_sessions.set_inactive(user, 'time')
                            await self.bot.send_message(chat_id=user,
                                                            text='Привет! Твоя подписка закончилась. Можешь продолжить общение с нейросетью, купив подписку')
                        except Exception as e:
                            if 'blocked' in str(e):
                                await self.db.set_blocked_user(user)
                                count_error_users_for_inactive += 1
                                pass
                            else:
                                print(e)
                                pass

                users = await self.db.get_sub_ending_users()
                k1 = str(len(users))
                k1_errors = 0
                k1_error_messages = []
                for user in users:


                        try:
                            await self.bot.send_message(chat_id=user,
                                                        text='Привет, напомиаем, что сегодня последний день действия подписки.')
                        except Exception as e:
                            await self.db.set_blocked_user(user)
                            k1_errors += 1
                            k1_error_messages.append(str(e))

                            pass

                users = await self.db.get_active_trial_users()

                k2 = str(len(users))
                count_error = 0
                count_send = 0
                error_messages = []
                for user in users:

                    try:

                        date = str(await self.db.get_end_time(user))[0:10]
                        date = date[8:10] + '.' + date[5:7] + '.' + date[0:4]
                        await self.bot.send_message(chat_id=user,
                                                    text='История чата автоматически сброшена. Токены обнулены, а ты можешь продолжить задавать вопросы GPT' + '\n' +
                                                    'Ваша подписка : ' + await self.db.get_sub_name_from_user(user) + '\n' + 'Закончится: ' +
                     date)
                        count_send += 1

                    except Exception as e:
                        count_error += 1
                        error_messages.append(str(e))
                        await self.db.set_blocked_user(user)
                        pass
                unique_error_messages = (set(error_messages))
                k1_error_messages = (set(k1_error_messages))

                try:
                    count_new_users = str( await self.db.count_new_users_trial())
                    count_sold = str(await self.db.count_new_users_not_trial())
                except:
                    count_new_users = '0'
                    count_sold = '0'
                    pass
                for admin_id in admin:
                    try:
                        await self.bot.send_message(chat_id=admin_id,
                                                    text='Отправили уведомление о конце заканчивающейся' + '\n' + 'Количество пользователей: ' + k1 + '\n'+
                                                    'Количество пользователей с ошибкой: ' + str(k1_errors) + '\n' + 'Уникальные ошибки: ' + str(k1_error_messages))
                        await self.bot.send_message(chat_id=admin_id,
                                                    text='Отправили уведомление о сбросе истории чата' + '\n' + 'Количество пользователей подошло: ' + k2 + '\n'+
                                                    'Количество пользователей с ошибкой: ' + str(count_error) + '\n' + 'Уникальные ошибки: ' + str(unique_error_messages))

                        await self.bot.send_message(chat_id=admin_id,
                                                    text='уведомление о дактивации подписки' + '\n' + 'Количество пользователей подошло: ' + str(count_users_for_inactive) + '\n'+
                                                    'Количество пользователей с ошибкой: ' + str(count_error_users_for_inactive))




                    except:
                        print('error in send notif to admin')
                        pass



            except Exception as e:


                await self.send_to_admin('error in clean history' + '\n' + str(e))




    async def admin(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        if not await  self.db.is_admin(update.message.from_user.id):
            return
        if self.openai.type_admin == 'personal':
            self.openai.type_admin = 'work'
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Режим админа изменен на рабочий',
            )
        else:
            self.openai.type_admin = 'personal'
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Режим админа изменен на личный',
            )




    async def buy(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        try:
            await self.add_offline(update.message.from_user.id, 'commandbuy')
        except Exception as e:
            logging.error(e)
            pass

        user_id = update.message.from_user.id

        if not await self.db.user_exists(user_id):
            logging.info(f'User {update.message.from_user.name} (id: {user_id}) is not in the database. Buying a subscription...')
            try:
                await update.message.reply_text(
                    message_thread_id=get_thread_id(update),
                    text='Введите команду /start для начала использования бота',
                )

            except:

                pass

        plan = await self.db.get_sub_name_from_user(user_id)

        if await  self.db.get_status(user_id) == 'active' and plan != 'trial':
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='''У вас уже есть активная подписка.

При повторной покупке неиспользованные токены не переносятся''',
            )
        discount = False
        prices = [80, 150, 260, 580]
        prices_new = prices
        prices_old = ['' for i in prices]
        # try:
        #     if await self.db.get_promo_used(user_id) == 0:
        #
        #         user_channel_status = await self.bot.get_chat_member(chat_id='@echokosmosa', user_id=user_id)
        #         logging.info(f'User {update.message.from_user.name} (id: {user_id}) status in channel: {user_channel_status.status}')
        #         if user_channel_status.status != 'left':
        #             discount = True
        #             prices_old = prices
        #             prices_new = [int(i * 0.9) for i in prices]
        #
        #         else:
        #             pass
        #
        #
        # except:
        #     pass
        subs = await  self.db.get_subs_for_sale()

        reply_markup_buttons = []

        try:
            for sub in subs:
                if discount:
                    button_text = sub['sub_name'] + ' ' + str(int(sub['price']*0.9)) + ' руб'
                else:
                    button_text = sub['sub_name'] + ' ' + str(sub['price']) + ' руб'
                button_callback = sub['sub_id']
                reply_markup_buttons.append([InlineKeyboardButton(text=button_text, callback_data=button_callback)])
        except Exception as e:
            print(traceback.format_exc())
            print(e)


        # Создаем разметку с кнопками из списка reply_markup_buttons
        reply_markup = InlineKeyboardMarkup(reply_markup_buttons)

        text = f'''
🗒Описание подписок:

<b>🔸GPT Mini🔸</b>
<blockquote>
💰  <i>{prices_old[0]}</i>  <b>{prices_new[0]} руб / 30 дней</b>

⚙️     <b>GPT-4-mini </b>: 500 тыс.

🔹 Чтение файлов
❌ Анализ и генерация фото

</blockquote>

<b>🔸Multi Light🔸</b>
<blockquote>
💰  <i>{prices_old[1]}</i>  <b>{prices_new[1]} руб / 30 дней</b>

⚙️  Доступно токенов при использовании: 

    <b>GPT-4      </b>: 100 тыс.
    <i>или</i>
    <b>GPT-4-mini  </b>: 1 млн.

🔹 Анализ фото
🔹 Генерация до 20 изображений
🔹 Чтение файлов 

</blockquote> 

<b>🔸Multi Standart🔸</b>
<blockquote>
💰  <i>{prices_old[2]}</i>  <b>{prices_new[2]} руб / 30 дней</b>

⚙️  Доступно токенов при использовании:

    <b>GPT-4      </b>: 200 тыс.
    <i>или</i>
    <b>GPT-4-mini  </b>: 2 млн.

🔹 Анализ фото и файлов
🔹 Генерация до 40 изображений
🔹 Чтение файлов 

</blockquote> 

<b>🔸Multi PRO🔸</b>
<blockquote>
💰  <i>{prices_old[3]}</i>  <b>{prices_new[3]} руб / 60 дней</b>

⚙️ Доступно токенов при использовании: 

    <b>GPT-4      </b>: 500 тыс.
    <i>или</i>
    <b>GPT-4-mini  </b>: 5 млн.

🔹 Анализ фото
🔹 Генерация до 100 изображений 
🔹 Чтение файлов 

</blockquote>

⚙️ Менять модель командой /model

📢  Переключив модель с GPT-4 на GPT-4-mini расход токенов снизится в 10 раз

<b>✨ Сравнение моделей:</b>
<blockquote>
GPT-4o         89%
GPT-4o-mini     82%

% — доля правильных отетов

</blockquote>

—————
<b>🔻Важно</b>
<blockquote>

❗️ Приблизительно 1000 токенов – 300 слов или 2300 символов или 1.5 стр. А4.

🔹 1 фото для анализа весит 1500 токенов
🔹 Генерация 1 изображения стоит от 5000 токенов 
🔹 <b>Поддерживаются форматы файлов: pdf, pptx, xlsx, txt, docx</b> – бот читает текст из файлов и добавляет его к запросу, обращаейте внимание на трату токенов

</blockquote>

Подробнее на сайте: brainstormai.ru

Выберите подписку или введите /cancel для отмены
'''


        await update.effective_message.reply_text(
            message_thread_id=get_thread_id(update),

            text=text,
            parse_mode = 'HTML',
            reply_markup=reply_markup
        )

        # try:
        #     if await self.db.get_promo_used(user_id) == 0:
        #         if not discount:
        #
        #             await update.effective_message.reply_text(
        #                 message_thread_id=get_thread_id(update),
        #                 parse_mode='HTML',
        #                 text='Подпишись на канал <a href="https://t.me/+lvsQbyECDwE0MDdi">@echokosmosa</a> и получи скидку 10%',
        #             )
        #         else:
        #             await update.effective_message.reply_text(
        #                 message_thread_id=get_thread_id(update),
        #                 parse_mode='HTML',
        #                 text='Скидка применена',
        #             )
        # except:
        #     pass
    async def send_end_of_subscription_message(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Ваша подписка закончилась, купите подписку',
        )
        await self.buy(update, context)

    async def temperature(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        user_id = update.message.from_user.id
        if not await self.db.get_edit_temp(user_id):
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка не позволяет изменять температуру',

            )

            return
        if not await self.is_active(update, context, update.message.from_user.id):
            await self.send_end_of_subscription_message(update, context)
            return

        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Введите температуру от 0 до 1.25',
        )
        user_id = update.message.from_user.id
        self.status[user_id] = 'set_temperature'

    async def role(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        if not await self.db.get_edit_role(update.message.from_user.id):
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Ваша подписка не позволяет изменять роль',
            )

            return
        if not await self.is_active(update, context, update.message.from_user.id):
            await self.send_end_of_subscription_message(update, context)
            return

        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='''Придумайте роль или введите /cancel для отмены

После ввода роли история чата автоматически сбросится''',
        )

        user_id = update.message.from_user.id
        self.status[user_id] = 'set_role'

    async def set_role(self, update: Update, context: ContextTypes.DEFAULT_TYPE, role):

        await update.message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Роль изменена',
        )
        user_id = update.message.from_user.id
        self.status[user_id] = 'prompt'
        self.openai.reset_chat_history(chat_id=update.effective_chat.id, content=message_text(update.message))
        self.openai.add_role_to_history(chat_id=update.effective_chat.id, content=role)

        sub_name = await self.db.get_sub_name_from_user(user_id)
        try:
            if sub_name =='ultimate admin':
                pass
            else:
                await self.db_analytics_for_sessions.role_edited(user_id)
        except Exception as e:
            logging.error(f'Error in set_role: {e}')
            pass
        # await self.db_analytics_for_month.add_role_edited(await self.db.get_sub_type(user_id))

    async def set_temperature(self, update: Update, context: ContextTypes.DEFAULT_TYPE, temperature):
        try:
            if float(temperature) <= 1.25 and float(temperature) >= 0.0:

                await update.message.reply_text(
                    message_thread_id=get_thread_id(update),
                    text='Температура изменена',
                )
                await self.db.set_custom_temp(update.message.from_user.id, temperature)
                user_id = update.message.from_user.id
                self.status[user_id] = 'prompt'
                sub_name = await self.db.get_sub_name_from_user(user_id)
                try:
                    if sub_name == 'ultimate admin':
                        pass
                    else:

                        await self.db_analytics_for_sessions.temp_edited(user_id)
                except Exception as e:
                    logging.error(f'Error in set_temperature: {e}')
                    pass

                # await self.db_analytics_for_month.add_temp_edited(await self.db.get_sub_type(user_id))

        except Exception as e:
            logging.error(f'Error in set_temperature: {e}')
            await update.message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Введите температуру от 0 до 1.25 или /cancel для отмены',
            )


    async def button(self, update: Update, context: ContextTypes.DEFAULT_TYPE):
        user_id = update.callback_query.from_user.id

        if update.callback_query.data == 'st-1':
            self.quality_list[user_id] = 'st-1'
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Качество изменено. Чтобы сгенерировать фото, введите  /imagine',
            )
            return
        elif update.callback_query.data == 'st-2':
            self.quality_list[user_id] = 'st-2'
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Качество изменено. Чтобы сгенерировать фото, введите  /imagine',
            )
            return
        elif update.callback_query.data == 'st-3':
            self.quality_list[user_id] = 'st-3'
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Качество изменено. Чтобы сгенерировать фото, введите  /imagine',
            )
            return
        elif update.callback_query.data == 'hd-1':
            self.quality_list[user_id] = 'hd-1'
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Качество изменено. Чтобы сгенерировать фото, введите  /imagine',
            )
            return
        elif update.callback_query.data == 'hd-2':
            self.quality_list[user_id] = 'hd-2'
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Качество изменено. Чтобы сгенерировать фото, введите  /imagine',
            )
            return
        elif update.callback_query.data == 'hd-3':
            self.quality_list[user_id] = 'hd-3'
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text='Качество изменено. Чтобы сгенерировать фото, введите  /imagine',
            )
            return
        elif await self.db.get_email(update.callback_query.from_user.id) == None:
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                parse_mode = 'HTML',
                text='''Пожалуйста, <b>введите email</b>, на который будет выслан чек или /cancel для отмены. 
                
С политикой кофиденциальности можно ознакомится на сайте https://brainstormai.ru/privacy-policy''',

            )

            self.status[user_id] = 'set_email'
            return
        await update.effective_message.reply_text(
            message_thread_id=get_thread_id(update),
            text='Ожидайте ссылку на оплату',
        )
        query = update.callback_query
        await query.answer()

        Configuration.account_id = self.config['shop_id']
        Configuration.secret_key = self.config['yookassa_key']
        discount = False
        # try:
        #     if await self.db.get_promo_used(user_id) == 0:
        #
        #
        #         user_channel_status = await self.bot.get_chat_member(chat_id='@echokosmosa', user_id=user_id)
        #         print(user_channel_status.status)
        #         if user_channel_status.status != 'left':
        #             discount = True
        #         else:
        #             pass
        # except Exception as e:
        #
        #     pass


        price = await self.db.get_price(query.data)

        if discount:
            price = price * 0.9





        sub_name = await self.db.get_sub_name(query.data)
        email = await self.db.get_email(query.from_user.id)
        try:
            payment_details =  payment.payment(price, sub_name, email)
        except Exception as e:
            logging.error(f'Error in create payment: {e}')
            if 'email' in str(e):

                await self.db.reset_email(query.from_user.id)
                await update.effective_message.reply_text(
                    message_thread_id=get_thread_id(update),
                    parse_mode = 'HTML',
                    text='''Некорректный email, пожалуйста, <b>введите email</b>, на который будет выслан чек.'''
                )
                self.status[user_id] = 'set_email'
                return

            else:
                await self.send_to_admin( 'error in create payment' + '\n' + str(e))
                await update.effective_message.reply_text(
                    message_thread_id=get_thread_id(update),
                    text='Ошибка при создании платежа, обратитесь в поддержку',
                )

                return


        await update.effective_message.reply_text(
            (payment_details['confirmation'])['confirmation_url'])
        payment_success = False

        try:
            payment_success = await payment.check_payment(payment_details['id'])
        except Exception as e:
            logging.error(f'Error in check payment: {e}')
            await self.send_to_admin( 'error in check payment' + '\n' + 'for user:' + str(user_id) + '\n' + str(e))



        if payment_success:

            try:
                await self.db.add_promo_used(user_id)
            except:
                pass


            user_id = update.callback_query.from_user.id
            sub_id = query.data
            order_id = payment_details['id']
            try:
                await self.activate_sub(user_id, query.data, order_id)

            except Exception as e:
                logging.error(f'Error in activate sub: {e}')
                await self.send_to_admin( 'error in activate sub' + '\n' + str(e))
                pass
            try:


                await self.send_to_admin('Платеж прошел' + '\n' + 'Пользователь: ' + str(user_id) + '\n' + 'Подписка: ' + sub_name + '\n' + 'Цена: ' + str(price) + '\n' + 'Email: ' + email)
            except Exception as e:
                logging.error(f'Error in send message: {e}')
                pass
            try:
                if await self.db.get_sub_multimodel(sub_id):
                    await self.db.set_user_model(user_id, 'gpt-4o-mini')
                    await update.effective_message.reply_text(
                        message_thread_id=get_thread_id(update),
                        text='Сейчас вы используете модель GPT 4 mini, расход токенов уменьшен в 5 раз, для смены модели на GPT-4 введите /model',
                    )
                await update.effective_message.reply_text("Платеж прошел")
            except Exception as e:

                pass

            try:
                income = price
                await self.db_analytics_for_sessions.new_sub_stats(user_id, sub_id, order_id, income)
                order_id = await self.db_analytics_for_sessions.get_sub_stats_id(user_id)
                order_info = await self.db.get_sub_info(sub_id)
                cost = order_info['cost']
                product = order_info['sub_name']
                income = order_info['price']
                client_id = await self.db.get_client_id(user_id)
                count = await self.db_analytics_for_sessions.count_orders(user_id)
                if count == 2 and client_id is not None:

                    await self.add_client(update, context, user_id, client_id)
                    await self.add_order(user_id, income, cost, order_id, product)
            except Exception as e:
                logging.error(f'Error in add order to metrika: {e}')
                await self.send_to_admin( 'error in add order metrika' + '\n' + str(e))
                pass


            # await self.db_analytics_for_month.add_income(sub_id, await self.db.get_price(sub_id))
            # await self.db_analytics_for_month.add_sold(sub_id)

        else:
            await update.effective_message.reply_text(
                message_thread_id=get_thread_id(update),
                text="Платеж не прошел, попробуйте еще раз"
            )

    async def activate_sub(self, user_id, sub_id, order_id_payment):

        try:

            await self.db_analytics_for_sessions.set_inactive(user_id, 'new_sub')



        except Exception as e:
            logging.error(f'Error in activate sub analytics1: {e}')
            await self.send_to_admin( 'error in activate sub analytics' + '\n' + str(e))

            pass
        self.openai.reset_chat_history(chat_id=user_id)


        await self.db.update_user(user_id, sub_id)
        try:

            income = await self.db.get_price(sub_id)
            cost = 0
            # await self.db_analytics_for_sessions.new_sub_stats(user_id, sub_id, order_id_payment, income)



        except Exception as e:
            logging.error(f'Error in activate sub analytics2: {e}')
            await self.send_to_admin( 'error in activate sub analytics2' + '\n' + str(e))
            pass


    async def is_in_time(self, update: Update, context: ContextTypes.DEFAULT_TYPE, user_id) -> bool:

        if await self.db.get_sub_type(user_id) == 2:
            return True
        elif await self.db.get_end_time(user_id) <= datetime.now().date():
            return False
        else:
            return True

    async def is_active(self, update: Update, context: ContextTypes.DEFAULT_TYPE, user_id: int) -> bool:
        if await self.db.get_sub_type(user_id) == 2:
            return True
        elif await self.db.get_status(user_id) == 'inactive':
            return False
        elif await self.is_in_time(update, context, user_id) == False:

            await self.db.set_inactive(user_id)



            self.openai.reset_chat_history(chat_id=user_id)
            sub_name = await self.db.get_sub_name_from_user(user_id)

            try:
                if sub_name == 'ultimate admin':
                    pass
                else:

                    await self.db_analytics_for_sessions.set_inactive(user_id, 'time')

            except Exception as e:
                print(traceback.format_exc())
                await self.send_to_admin( 'error in is active' + '\n' + str(e))
                pass

            # await self.db_analytics_for_month.add_expired(sub_type)
            # await self.db_analytics_for_month.add_expired_time(sub_type)

            return False

        else:
            return True

    def is_email(self, email):


        if re.match(r"[^@]+@[^@]+\.[^@]+", email) and len(email) > 5:
            return True
        else:
            return False

    async def is_in_tokens(self, update: Update, context: ContextTypes.DEFAULT_TYPE, user_id) -> bool:
        if await self.db.get_sub_type(user_id) == 2:
            return True
        elif await self.db.get_used_tokens(user_id) >= await self.db.get_max_tokens(user_id):
            return False
        else:
            return True

    async def is_input_in_tokens(self, update: Update, context: ContextTypes.DEFAULT_TYPE, user_id,
                                 tokens_input, model_config=None) -> bool:

        if model_config is None:
            model_config = {'multimodel_3': False, 'multi_k': 1}
        multimodel_3 = True
        multi_k = model_config['multi_k']

        if await self.db.get_sub_type(user_id) == 2:
            return True
        elif multimodel_3 and await self.db.get_max_tokens(user_id)*multi_k - (await self.db.get_used_tokens(user_id)*multi_k + tokens_input) <= 10:

                return False
        elif not multimodel_3 and  await self.db.get_max_tokens(user_id) - (await self.db.get_used_tokens(user_id) + tokens_input) <= 10:

            return False
        else:
            return True

    async def calc_end_time(self, user_id):
        current_date = datetime.now().date()
        end_date = current_date + timedelta(days=await self.db.get_duration_sub(user_id))
        await self.db.set_end_time(user_id, end_date)

    async def inactivate(self, user_id, update: Update, context: ContextTypes.DEFAULT_TYPE):


            plan = await self.db.get_sub_name_from_user(user_id)
            if plan == 'trial':
                await update.effective_message.reply_text(
                    message_thread_id=get_thread_id(update),
                    text='Ваш лимит токенов на сегодня закончился, купите подписку или подождите до завтра',
                )

            else:
                await update.effective_message.reply_text(
                    message_thread_id=get_thread_id(update),
                    text='Ваш лимит токенов закончился, купите подписку',
                )

                await self.db.set_inactive(user_id)
                self.openai.reset_chat_history(chat_id=user_id)

                try:
                    sub_name = await self.db.get_sub_name_from_user(user_id)
                    if sub_name == 'ultimate admin':
                        pass
                    else:

                        await self.db_analytics_for_sessions.set_inactive(user_id, 'tokens')

                except Exception as e:
                    logging.error(f'Error in inactivate: {e}')
                    await self.send_to_admin('error in is active' + '\n' + str(e))
                    pass

                # await  self.db_analytics_for_month.add_expired(plan)
                # await self.db_analytics_for_month.add_expired_tokens(plan)

                await self.buy(update, context)



    async def prompt(self, update: Update, context: ContextTypes.DEFAULT_TYPE):

        try:
            """
            React to incoming messages and respond accordingly.
            """
            if update.edited_message or not update.message or update.message.via_bot:
                return

            # if not await self.check_allowed_and_within_budget(update, context):
            #     return

            logging.info(
                f'New message received from user {update.message.from_user.name} (id: {update.message.from_user.id})'
                )


            chat_id = update.effective_chat.id
            user_id = update.message.from_user.id


            if not await self.db.user_exists(user_id):
                logging.error(f'User {update.message.from_user.name} (id: {user_id}) is not in the database. In prompt')
                try:
                    await update.message.reply_text(
                        message_thread_id=get_thread_id(update),
                        text='Введите команду /start для начала использования бота',
                    )
                except:

                    pass

            model_config = await self.db.get_model_config(user_id)




            photo_list = []
            file_text = ""

            # Handle photos
            if update.message.photo:
                try:
                    await self.db_analytics_for_sessions.photo_send(user_id)
                    logging.info(f'Photo send: {user_id}')
                except:
                    pass
                model_name = await self.db.get_user_model(user_id)

                if  model_name != 'gpt-4':

                    await update.message.reply_text(
                        message_thread_id=get_thread_id(update),
                        text='Анализировать фото может только gpt-4. Чтобы сменить модель — введите команду /model',
                    )
                    try:
                        self.prompts[chat_id] -= 1
                    except:
                        pass
                    return

                file = update.message.photo[-1].file_id  # get the file_id of the largest size photo
                obj = await context.bot.get_file(file)
                out = io.BytesIO()
                await obj.download_to_memory(out=out)
                out.seek(0)  # reset file pointer to the beginning
                base64_image = base64.b64encode(out.read()).decode('utf-8')

                # Add the base64 image to the list
                photo_list.append(base64_image)

            # Handle document files
            if update.message.document:
                try:
                   file = await context.bot.get_file(update.message.document.file_id)
                except Exception as e:
                    print(e)
                    if 'File is too big' in str(e):
                        await update.message.reply_text(
                            message_thread_id=get_thread_id(update),
                            text='Файл слишком большой, попробуйте загрузить файл меньшего размера'
                        )
                    else:
                        await update.message.reply_text(
                            message_thread_id=get_thread_id(update),
                            text='Произошла ошибка, попробуйте еще раз'
                        )
                    return



                
                file_obj = io.BytesIO()
                await file.download_to_memory(out=file_obj)
                file_obj.seek(0)

                mime_type = update.message.document.mime_type
                file_name = update.message.document.file_name

                if mime_type == 'application/pdf':
                    pdf_reader = PyPDF2.PdfReader(file_obj)
                    for page in pdf_reader.pages:
                        file_text += page.extract_text() + "\n"

                elif mime_type == 'text/plain':
                    file_text = file_obj.read().decode('utf-8')

                elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    doc = docx.Document(file_obj)
                    file_text = "\n".join([paragraph.text for paragraph in doc.paragraphs])

                elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    prs = Presentation(file_obj)
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                file_text += shape.text + "\n"

                elif mime_type == 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet':
                    wb = openpyxl.load_workbook(file_obj, data_only=True)
                    file_text = "Excel file content:\n\n"
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        file_text += f"Sheet: {sheet}\n"
                        
                        # Get column headers
                        headers = [cell.value for cell in ws[1] if cell.value is not None]
                        
                        # Calculate column widths
                        col_widths = [max(len(str(cell.value)) if cell.value is not None else 0 for cell in col) for col in ws.columns]
                        
                        # Add headers
                        header_row = "| " + " | ".join(f"{str(headers[i]):<{col_widths[i]}}" for i in range(len(headers))) + " |"
                        file_text += header_row + "\n"
                        file_text += "|" + "|".join("-" * (width + 2) for width in col_widths) + "|\n"
                        
                        # Add data rows
                        for row in ws.iter_rows(min_row=2, values_only=True):
                            data_row = "| " + " | ".join(f"{str(row[i]):<{col_widths[i]}}" if row[i] is not None else " "*col_widths[i] for i in range(len(row))) + " |"
                            file_text += data_row + "\n"
                        
                        file_text += "\n"  # Add space between sheets

                        

                     


                        


                else:
                    await update.message.reply_text(
                        message_thread_id=get_thread_id(update),
                        text="Неподдерживаемый формат файла. Поддерживаемые форматы: PDF, TXT, DOCX, PPTX, XLSX."
                    )
                    return

            # Create the prompt
            if photo_list:
                prompt = []
                # Add a text prompt for photos
                if update.message.caption:
                    text = update.message.caption
                else:
                    text = 'Analyze the attached image(s)'
                
                prompt.append({
                    "type": "text",
                    "text": text
                })

                # Add image prompts for each photo in the list
                for base64_image in photo_list:
                    prompt.append({
                        "type": "image_url",
                        "image_url": {"url": f"data:image/jpeg;base64,{base64_image}"}
                    })
            elif file_text:
                # For document files, use the extracted text as the prompt
                if update.message.caption:
                    text = update.message.caption
                else:
                    text = f'Прочитай и запомни информацию {file_name}, попроси пользователя сформулировать вопрос по ней'
                prompt = f"{text}\n\nFile content:\n{file_text}..."  # Limit to first 2000 characters

                model = await self.db.get_user_model(user_id)

                tokens =  self.openai.count_tokens(([{"role": "user", "content": prompt}]), model)
                tokens = tokens//model_config['multi_k']

                self.openai.add_to_history(chat_id, "user", prompt)

                await update.message.reply_text(
                    message_thread_id=get_thread_id(update),
                    text=f'''Вы прикрепили файл: чтение этого файла стоит ~{tokens} ток. — для модели gpt-4-mini

Чтобы отменить — введите /reset

Чтобы продолжить — просто напишите вопрос, который хотели задать'''

                )
                return

            else:
                prompt = update.message.text

            if user_id not in self.status:
                self.status[user_id] = 'prompt'

            if update.message:
                if self.status[user_id] == 'set_email':
                    if self.is_email(update.message.text) == False:
                        await update.message.reply_text(
                            message_thread_id=get_thread_id(update),
                            text='Неверный формат email',
                        )
                        return
                    await self.db.set_email(user_id, update.message.text)
                    await update.message.reply_text(
                        message_thread_id=get_thread_id(update),
                        text='Email установлен, выберите подписку',
                    )
                    self.status[user_id] = 'prompt'

                    await self.buy(update, context)
                    return

        



                elif self.status[user_id] == 'image':

                    self.status[user_id] = 'prompt'

                    prompt = update.message.text
                    res = self.get_quality(user_id)
                    quality = res['quality']
                    size = res['size']
                    price = res['price']
                    tokens = res['tokens']




                    url = await self.openai.generate_image( quality, size, prompt)
                    if url == False:
                        await update.message.reply_text(
                            message_thread_id=get_thread_id(update),
                            text='Произошла ошибка, попробуйте еще раз',
                        )
                        return
                    else:

                        try:
                            await update.message.reply_photo(
                            photo=url,
                            message_thread_id=get_thread_id(update),
                            )



                            await self.db.update_used_tokens(user_id, tokens)
                        except Exception as e:
                            if 'Failed to get http url content' in e:
                                print(str(url))

                                await update.message.reply_text(
                                    text=str(url),
                                )
                            elif 'safety system' in e:
                                await update.message.reply_text(
                                    message_thread_id=get_thread_id(update),
                                    text='Ваш запрос не прошел модерацию',
                                )

                            else:
                                await update.message.reply_text(
                                    message_thread_id=get_thread_id(update),
                                    text='Произошла ошибка, попробуйте еще раз',
                                )
                            return




                        try:

                            sub_name = await self.db.get_sub_name_from_user(user_id)
                            cost={ 'output' : price,'input' : price}
                            if sub_name == 'ultimate admin':
                                await self.db_admin_stats.add(user_id,1/2,1/2, cost)
                                print('admin')
                            else:
                                print('not admin')
                                await self.db_analytics_for_sessions.image_generated(user_id)
                                await self.db_statistic_by_day.add(user_id, 1/2, 1/2, cost)
                        except Exception as e:
                            print('error in analytics image' + '\n' + str(e))
                            pass

                        return

                elif self.status[user_id] == 'admin_message':

                    self.status[user_id] = 'prompt'
                    users = await self.db.get_all_users()
                    print(users)
                    k = 0
                    err =0
                    for user in users:
                        try:
                            await self.bot.send_message(chat_id=user, text=update.message.text,parse_mode = 'HTML')
                            k+=1
                        except Exception as e:
                         
                            err+=1
                            await self.db.set_blocked_user(user)

                            pass

                    print(k)
                    await self.send_to_admin('send message to all users' + '\n' + str(k) + 'error: ' + str(err))
                    return

                elif self.status[user_id] == 'admin_message_to_act':

                    self.status[user_id] = 'prompt'
                    users = await self.db.get_act_users()
                    k = 0
                    err =0
                    for user in users:
                        try:
                            await self.bot.send_message(chat_id=user, text=update.message.text,parse_mode = 'HTML')
                            k+=1
                        except Exception as e:
                            err+=1
                            await self.db.set_blocked_user(user)

                            pass

                    print(k)
                    await self.send_to_admin('send message to all users' + '\n' + str(k) + 'error: ' + str(err))
                    return

                elif not await self.is_active(update, context, user_id):
                    await self.send_end_of_subscription_message(update, context)
                    return

                if self.status[user_id] == 'set_role':
                    await self.set_role(update, context, update.message.text)

                    return
                elif self.status[user_id] == 'set_temperature':

                    await self.set_temperature(update, context, update.message.text)





                elif self.status[user_id] == 'prompt':


                    plan = await self.db.get_sub_type(user_id)
                    plan_name = await self.db.get_sub_name_from_user(user_id)
                    model = await self.db.get_model(user_id)
                    date = datetime.now().date()
                    last_message = await self.db.get_last_message(user_id)



                    self.prompts[chat_id] = self.prompts.get(chat_id, 0) + 1
                    print(self.prompts[chat_id])

                    if self.prompts[chat_id] >1:
                        available_tokens = await self.db.get_max_tokens(user_id) - await self.db.get_used_tokens(user_id)
                        if (default_max_tokens(model) + self.openai.count_tokens(([{"role": "user", "content": prompt}]), model) ) * self.prompts[chat_id] > available_tokens:
                            await update.effective_message.reply_text(
                                message_thread_id=get_thread_id(update),
                                text='К сожалению у вас заканчиваются токены, поэтому мы не можем выполнить ваши запросы параллельно. Пожалуйста, подождите ответа на прошлое сообщение и попробуйте снова',
                            )
                            self.prompts[chat_id] -= 1
                            return

                    if last_message != date:



                        if plan_name == 'trial':
                            await self.db.set_used_tokens(user_id, 0)
                            self.openai.reset_chat_history(chat_id=user_id)
                        elif  last_message == date - timedelta(days=2):
                            self.openai.reset_chat_history(chat_id=user_id)

                        await self.db.set_last_message(user_id, date)


                    self.last_message[chat_id] = prompt



                    model_config = await self.db.get_model_config(update.effective_chat.id)

                    if model_config is None:
                        await update.effective_message.reply_text(
                            message_thread_id=get_thread_id(update),
                            text='Ваша подписка закончилась, купите подписку',
                        )
                        return

                    


                    tokens_in_message = self.openai.count_tokens(([{"role": "user", "content": str(prompt)}]), model_config['model'])
                    tokens_input = tokens_in_message + self.openai.get_conversation_stats(chat_id=chat_id, model=model_config['model'])[1]

                    while not await self.is_input_in_tokens(update, context, user_id, tokens_input, model_config):
                        try:
                            if self.openai.remove_messages(chat_id):
                                tokens_input = tokens_in_message + self.openai.get_conversation_stats(chat_id=chat_id, model=model_config['model'])[1]
                            else:

                                await update.effective_message.reply_text(
                                    message_thread_id=get_thread_id(update),
                                    text='Осталось ' + ' ' + str(
                                        await self.db.get_max_tokens(user_id) - await self.db.get_used_tokens(
                                            user_id)) + ' токенов' + '\n' + 'Ваше сообщение слишком длинное, сократите его ',
                                )
                                self.prompts[chat_id] -= 1
                                return
                        except Exception as e:
                            await self.send_to_admin('error in remove messages' + '\n' + str(e))
                            break



                    try:
                        total_tokens = 0
                        input_tokens = 0

                        if self.config['stream']:



                            # get photo from message and send to ai
                            # get photo from message and send to ai

                            text = update.message.text

                            base64_image = None
                            # get photo from message and send to ai



                            async def _reply():
                                nonlocal total_tokens
                                await update.effective_message.reply_chat_action(
                                    action=constants.ChatAction.TYPING,
                                    message_thread_id=get_thread_id(update)
                                )








                                stream_response = self.openai.get_chat_response_stream(chat_id=chat_id, query=prompt,
                                                                                           model_config=model_config, sub_type = plan)

                                i = 0
                                prev = ''
                                sent_message = None
                                backoff = 0
                                stream_chunk = 0

                                async for content, tokens in stream_response:

                                    if len(content.strip()) == 0:
                                        continue

                                    stream_chunks = split_into_chunks(content)
                                    if len(stream_chunks) > 1:
                                        content = stream_chunks[-1]
                                        if stream_chunk != len(stream_chunks) - 1:
                                            stream_chunk += 1
                                            try:
                                                await edit_message_with_retry(context, chat_id,
                                                                              str(sent_message.message_id),
                                                                              stream_chunks[-2])
                                            except:
                                                pass
                                            try:
                                                sent_message = await update.effective_message.reply_text(
                                                    message_thread_id=get_thread_id(update),
                                                    text=content if len(content) > 0 else "..."
                                                )
                                            except:
                                                pass
                                            continue

                                    cutoff = get_stream_cutoff_values(update, content)
                                    cutoff += backoff

                                    if i == 0:
                                        try:
                                            if sent_message is not None:
                                                await context.bot.delete_message(chat_id=sent_message.chat_id,
                                                                                 message_id=sent_message.message_id)
                                            sent_message = await update.effective_message.reply_text(
                                                message_thread_id=get_thread_id(update),
                                                reply_to_message_id=get_reply_to_message_id(self.config, update),
                                                text=content
                                            )
                                        except:
                                            continue

                                    elif abs(len(content) - len(prev)) > cutoff or tokens != 'not_finished':
                                        prev = content

                                        try:
                                            use_markdown = tokens != 'not_finished'
                                            await edit_message_with_retry(context, chat_id, str(sent_message.message_id),
                                                                          text=content, markdown=use_markdown)

                                        except RetryAfter as e:
                                            backoff += 5
                                            try:
                                                # error Flood control exceeded. Retry in 130 seconds. we need to find string with seconds after words 'Retry in' and convert it to int
                                                wait_time = int(re.search(r'\d+', e.message).group())



                                                await update.effective_message.reply_text(
                                                    message_thread_id=get_thread_id(update),
                                                    text=f'''Бот вернулся к вам. Он не мог отвечать из-за спам-оганичений Telegram-а

В среднем блок длится от 3 секунд до 3 минут. Новые запросы — не помогают. Нужно просто подождать'''
                                                )
                                            except:
                                                pass


                                            await asyncio.sleep(e.retry_after)

                                            continue

                                        except TimedOut:
                                            backoff += 5
                                            await asyncio.sleep(0.5)
                                            continue

                                        except Exception:
                                            backoff += 5

                                            continue
                                        await asyncio.sleep(0.01)



                                    i += 1
                                    if tokens != 'not_finished':
                                        total_tokens = int(tokens)


                            await wrap_with_indicator(update, context, _reply, constants.ChatAction.TYPING)
                            self.prompts[chat_id] -= 1
                            # await self.db.update_used_tokens(user_id, total_tokens)

                            # await self.db_analytics_for_month.add_input_tokens(plan, input_tokens)
                            # await self.db_analytics_for_month.add_total_tokens(plan, total_tokens)
                            # await self.db_analytics_for_month.add_output_tokens(plan, total_tokens - input_tokens)
                            # await self.db_analytics_for_periods.add(plan, total_tokens)





                            if await self.is_in_tokens(update, context, user_id) == False:

                                await self.inactivate(user_id, update, context)














                    except Exception as e:
                        # traceback
                        await self.send_to_admin('error in prompt' + '\n' + str(e) + str(traceback.format_exc()))

                        self.prompts[chat_id] = 0

                        logging.exception(traceback.format_exc())
                        logging.error(f'Error in prompt: {e}')
                        await update.effective_message.reply_text(
                            message_thread_id=get_thread_id(update),
                            text='Произошла ошибка, попробуйте еще раз',
                        )
        except Exception as e:
            logging.error(f'Error in prompt: {e}')
            logging.exception(traceback.format_exc())




    async def post_init(self, application: Application) -> None:
        """
        Post initialization hook for the bot.
        """
        await application.bot.set_my_commands(self.group_commands, scope=BotCommandScopeAllGroupChats())
        await application.bot.set_my_commands(self.commands)

    async def run_ptb(self):

        """
        Runs the bot indefinitely until the user presses Ctrl+C
        """
        application = ApplicationBuilder() \
            .token(self.config['token']) \
            .proxy_url(self.config['proxy']) \
            .get_updates_proxy_url(self.config['proxy']) \
            .post_init(self.post_init) \
            .concurrent_updates(True) \
            .build()

        application.add_handler(CommandHandler('reset', self.reset))
        application.add_handler(CommandHandler('help', self.help))
        # application.add_handler(CommandHandler('image', self.image))
        application.add_handler(CommandHandler('start', self.start))
        application.add_handler(CommandHandler('buy', self.buy))
        application.add_handler(CommandHandler('stats', self.stats))
        application.add_handler(CommandHandler('resend', self.resend))
        application.add_handler(CommandHandler('cancel', self.cancel))
        application.add_handler(CommandHandler('role', self.role))
        application.add_handler(CommandHandler('send_to_all', self.send_to_all))
        application.add_handler(CommandHandler('send_to_act', self.send_to_act))

        application.add_handler(CommandHandler('admin', self.admin))
        application.add_handler(CommandHandler('save', self.save))
        application.add_handler((CommandHandler('delete', self.delete)))
        application.add_handler(CommandHandler('model', self.model))
        application.add_handler(CommandHandler('imagine', self.imagine))
        application.add_handler(CommandHandler('fluxdev', self.fluxdev))
        application.add_handler(CommandHandler('fluxpro', self.fluxpro))
        application.add_handler(CommandHandler('quality', self.quality))

        application.add_handler(CommandHandler('orders', self.orders))
        application.add_handler(CommandHandler('poll', self.poll))
        application.add_handler(CommandHandler('test_poll', self.test_poll))
        application.add_handler(CommandHandler("result", self.show_results))


        application.add_handler(PollAnswerHandler(self.receive_poll_answer))


        application.add_handler(CommandHandler('temperature', self.temperature))

        application.add_handler(CallbackQueryHandler(self.button))
        application.add_handler(CommandHandler(
            'chat', self.prompt, filters=filters.ChatType.GROUP | filters.ChatType.SUPERGROUP)
        )
        application.add_handler(MessageHandler( (~filters.COMMAND), self.prompt))
        application.add_error_handler(error_handler)

        # application.run_polling()
        await application.initialize()
        await application.start()
        await application.updater.start_polling()
        try:
            while application.running:
                await asyncio.sleep(1)
            logging.info("App stopped running for some other reason, shutting down...")
        except asyncio.CancelledError:
            logging.info('Bot stopped')
            await application.stop()
        finally:
            await application.stop()
            await application.shutdown()

    async def run_other(self, name):
        try:
            # do other async stuff, just sleeping here
            # await self.send_notif()
            await self.send_notif()
            logging.info(f"Other {name} stopped")
        except asyncio.CancelledError:
            logging.info(f"Other {name} stopped")
            # clean up here

    async def main(self):
            #linux
            tasks = [
                asyncio.create_task(self.run_ptb()),
                asyncio.create_task(self.run_other("send_notif"))
            ]
            try:

                for sig in (signal.SIGINT, signal.SIGTERM):
                    loop = asyncio.get_event_loop()
                    loop.add_signal_handler(sig, lambda sig=sig: asyncio.create_task(self.shutdown(tasks)))
                asyncio.get_event_loop() \
                    .add_signal_handler(signal.SIGTERM,
                                        lambda: asyncio.create_task(self.shutdown(tasks)))

                await asyncio.gather(*tasks)
            except:
                print('windows')
            # windows
                try:

                    await asyncio.gather(*tasks)
                except KeyboardInterrupt:
                    logging.info("Received Ctrl+C, shutting down gracefully.")
                    await self.shutdown(tasks)


    async def shutdown(self,tasks):
        logging.log.warning("Received SIGTERM, cancelling tasks...")
        for task in tasks:
            task.cancel()

    def run(self):
        try:
            asyncio.run(self.main())
        except Exception as e:
            print(e)


